import os
import docx
import pptx
import fitz  # PyMuPDF
import pandas as pd
from rich.console import Console
from rich.markdown import Markdown

console = Console()


def read_pptx(filepath):
    """Read and return PPTX file content."""
    try:
        presentation = pptx.Presentation(filepath)
        text = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)  # Return all text as a string
    except Exception as e:
        pass


def read_docx(filepath):
    """Read and return DOCX file content."""
    try:
        doc = docx.Document(filepath)
        return "\n".join(
            paragraph.text for paragraph in doc.paragraphs
        )  # Return all paragraphs as a string
    except Exception as e:
        pass


def read_pdf(filepath):
    """Read and return PDF file content."""
    try:
        doc = fitz.open(filepath)
        return "\n".join(
            page.get_text("text") for page in doc
        )  # Return all pages as a string
    except Exception as e:
        pass


def read_excel(filepath):
    """Read and return Excel file content."""
    try:
        df = pd.read_excel(filepath)
        return df.to_string(index=False)  # Return DataFrame as string
    except Exception as e:
        pass


def read_txt(filepath):
    """Read and return TXT file content."""
    try:
        with open(filepath, "r", encoding="utf-8") as file:
            return file.read()  # Return file content
    except Exception as e:
        pass


def upload_video(video_file_name, client):
    """
    Upload a video file to the GenAI server and return the uploaded file object.

    Parameters:
        video_file_name (str): The path to the video file to be uploaded.
        client (genai.Client): The GenAI client object to use for the upload.

    Returns:
        genai.File: The uploaded video file object.
    """
    video_file = client.files.upload(file=video_file_name)

    while video_file.state == "PROCESSING":
        console.print(Markdown("Video processing in progress..."), style="i cyan")
        video_file = client.files.get(name=video_file.name)

    if video_file.state == "FAILED":
        raise ValueError(video_file.state)
    console.print(
        Markdown(f"Video uploaded successfully! {video_file.uri}"), style="i green"
    )

    return video_file


def read_file(filepath):
    """Determine file type and call the appropriate function."""
    filepath = os.path.abspath(filepath)  # Convert relative to absolute path
    if not os.path.exists(filepath):
        pass

    ext = os.path.splitext(filepath)[1].lower()

    if ext == ".docx":
        return read_docx(filepath)
    elif ext == ".pptx":
        return read_pptx(filepath)
    elif ext == ".pdf":
        return read_pdf(filepath)
    elif ext in [".xls", ".xlsx"]:
        return read_excel(filepath)
    elif ext in [
        ".txt",
        ".py",
        ".json",
        ".html",
        ".css",
        ".js",
        ".md",
        ".c",
        ".cc",
        ".h",
        ".hpp",
        ".cs",
        ".cpp",
        ".java",
        ".php",
        ".sql",
        ".xml",
        ".yaml",
        ".yml",
        ".csv",
        ".tsv",
        ".log",
        ".bat",
        ".sh",
        ".ps1",
        ".psm1",
        ".psd1",
        ".ps1xml",
        ".pssc",
        # ".ini",
        # ".conf",
        # ".cfg",
        # ".env",
        # ... Add more file extensions as needed
    ]:
        return read_txt(filepath)
    else:
        pass


if __name__ == "__main__":
    file_path = input("Enter file path (absolute or relative): ").strip()
    try:
        print(read_file(file_path))
    except Exception as e:
        print("Error:", e)
